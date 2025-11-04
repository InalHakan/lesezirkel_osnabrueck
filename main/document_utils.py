import os
import tempfile
from io import BytesIO
from django.http import HttpResponse
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import PyPDF2


class DocumentConverter:
    """Utility class for converting various document formats to PDF"""
    
    @staticmethod
    def convert_to_pdf(file_path, original_filename):
        """
        Convert a document to PDF format
        Returns a BytesIO object containing the PDF data
        """
        file_extension = os.path.splitext(original_filename)[1].lower()
        
        if file_extension == '.pdf':
            # Already PDF, just return the file
            with open(file_path, 'rb') as f:
                return BytesIO(f.read())
        
        elif file_extension in ['.doc', '.docx']:
            return DocumentConverter._convert_docx_to_pdf(file_path)
        
        elif file_extension in ['.xls', '.xlsx']:
            return DocumentConverter._convert_excel_to_pdf(file_path)
        
        elif file_extension in ['.ppt', '.pptx']:
            return DocumentConverter._convert_pptx_to_pdf(file_path)
        
        elif file_extension in ['.txt', '.rtf']:
            return DocumentConverter._convert_text_to_pdf(file_path)
        
        else:
            # For unsupported formats, create a simple PDF with file info
            return DocumentConverter._create_info_pdf(original_filename)
    
    @staticmethod
    def _convert_docx_to_pdf(file_path):
        """Convert DOCX to PDF"""
        try:
            doc = DocxDocument(file_path)
            buffer = BytesIO()
            
            # Create PDF document
            pdf_doc = SimpleDocTemplate(buffer, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add document content
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    p = Paragraph(paragraph.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            return DocumentConverter._create_error_pdf(f"DOCX Conversion Error: {str(e)}")
    
    @staticmethod
    def _convert_excel_to_pdf(file_path):
        """Convert Excel to PDF"""
        try:
            wb = load_workbook(file_path)
            buffer = BytesIO()
            
            pdf_doc = SimpleDocTemplate(buffer, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # Add sheet title
                title = Paragraph(f"Sheet: {sheet_name}", styles['Heading1'])
                story.append(title)
                story.append(Spacer(1, 12))
                
                # Add sheet data
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        p = Paragraph(row_text, styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 24))
            
            pdf_doc.build(story)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            return DocumentConverter._create_error_pdf(f"Excel Conversion Error: {str(e)}")
    
    @staticmethod
    def _convert_pptx_to_pdf(file_path):
        """Convert PowerPoint to PDF"""
        try:
            prs = Presentation(file_path)
            buffer = BytesIO()
            
            pdf_doc = SimpleDocTemplate(buffer, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            for i, slide in enumerate(prs.slides):
                # Add slide title
                title = Paragraph(f"Slide {i + 1}", styles['Heading1'])
                story.append(title)
                story.append(Spacer(1, 12))
                
                # Extract text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        p = Paragraph(shape.text, styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 24))
            
            pdf_doc.build(story)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            return DocumentConverter._create_error_pdf(f"PowerPoint Conversion Error: {str(e)}")
    
    @staticmethod
    def _convert_text_to_pdf(file_path):
        """Convert text file to PDF"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            buffer = BytesIO()
            pdf_doc = SimpleDocTemplate(buffer, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Split content into paragraphs
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    p = Paragraph(para.strip(), styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            return DocumentConverter._create_error_pdf(f"Text Conversion Error: {str(e)}")
    
    @staticmethod
    def _create_info_pdf(filename):
        """Create an info PDF for unsupported formats"""
        buffer = BytesIO()
        pdf_doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        title = Paragraph("Document Information", styles['Title'])
        story.append(title)
        story.append(Spacer(1, 24))
        
        info = Paragraph(f"Original file: {filename}", styles['Normal'])
        story.append(info)
        story.append(Spacer(1, 12))
        
        message = Paragraph("This file format is not supported for PDF conversion. Please download the original file.", styles['Normal'])
        story.append(message)
        
        pdf_doc.build(story)
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def _create_error_pdf(error_message):
        """Create an error PDF"""
        buffer = BytesIO()
        pdf_doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        title = Paragraph("Conversion Error", styles['Title'])
        story.append(title)
        story.append(Spacer(1, 24))
        
        error = Paragraph(error_message, styles['Normal'])
        story.append(error)
        
        pdf_doc.build(story)
        buffer.seek(0)
        return buffer